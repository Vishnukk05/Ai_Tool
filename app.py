import os
import logging
import datetime
import uuid
import time
import base64
import psutil
from io import BytesIO

# --- FLASK & SERVER IMPORTS ---
from flask import Flask, render_template, request, jsonify, Response, session
from dotenv import load_dotenv

# --- AI & MEDIA IMPORTS ---
from groq import Groq
from gtts import gTTS
from xhtml2pdf import pisa
from pptx import Presentation
import speech_recognition as sr
import PIL.Image
from moviepy.video.io.VideoFileClip import VideoFileClip

# --- LOAD ENV ---
basedir = os.path.abspath(os.path.dirname(__file__))
env_path = os.path.join(basedir, '.env')
load_dotenv(env_path)

app = Flask(__name__)

# --- CRITICAL: SECRET KEY ---
app.secret_key = os.environ.get("SECRET_KEY", "super_secret_admin_key_12345") 

# --- LOGGING SETUP ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- CONFIGURATION ---
API_KEY = os.environ.get("GROQ_API_KEY") 
STATIC_FOLDER = os.path.join(basedir, 'static')
if not os.path.exists(STATIC_FOLDER): 
    os.makedirs(STATIC_FOLDER)

# --- STATS ---
global_stats = {
    "text_gen": 0, "audio_gen": 0, "transcribe": 0, "pdf_gen": 0, 
    "chat_msgs": 0, "code_review": 0, "quiz_gen": 0,
    "file_conv": 0, "compression": 0, "vid_audio": 0
}

def increment_stat(field_name):
    try:
        if field_name in global_stats: 
            global_stats[field_name] += 1
    except: pass

# --- ADMIN CREDENTIALS ---
ADMIN_USER = "Admin"
ADMIN_PASS = "M@nojkumarkk@2343"

# ==============================================================================
#                               HELPER FUNCTIONS
# ==============================================================================

def get_chat_response(messages):
    """
    Dedicated function for Chat that supports conversation history (Memory).
    """
    if not API_KEY:
        return "Error: API Key missing."

    try:
        client = Groq(api_key=API_KEY)
        
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile", 
            messages=messages,
            # Lower temperature = More factual/accurate.
            temperature=0.3, 
            max_tokens=1024,
            top_p=1,
        )
        return completion.choices[0].message.content
    except Exception as e:
        print(f"❌ Groq API Error: {e}")
        return "I apologize, but I am currently experiencing high traffic or an error. Please try again."

def get_safe_ai_response(prompt):
    """
    Wrapper for single-turn tasks (Quiz, Email, PPT, etc.)
    """
    if not API_KEY:
        print("❌ Error: GROQ_API_KEY not found in .env")
        return None

    try:
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile", 
            messages=[
                {"role": "system", "content": "You are a helpful AI assistant. Output clean text."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.6,
            max_tokens=2048,
        )
        return completion.choices[0].message.content
    except Exception as e:
        print(f"❌ Groq API Error: {e}")
        return None

# ==============================================================================
#                               CORE ROUTES
# ==============================================================================

@app.route('/')
def index(): 
    return render_template('index.html')

@app.route('/health')
def health_check():
    return "OK", 200

# --- AUTH ROUTES ---
@app.route('/login', methods=['POST'])
def login():
    data = request.json
    if data.get('username') == ADMIN_USER and data.get('password') == ADMIN_PASS:
        session['is_admin'] = True
        return jsonify({"success": True})
    return jsonify({"success": False, "error": "Invalid Credentials"}), 401

@app.route('/logout', methods=['POST'])
def logout():
    session.pop('is_admin', None)
    return jsonify({"success": True})

@app.route('/check-auth', methods=['GET'])
def check_auth():
    return jsonify({"is_admin": session.get('is_admin', False)})

# --- STATS ROUTES ---
@app.route('/api/stats')
def get_stats():
    if not session.get('is_admin'):
        return jsonify({"cpu": 0, "ram": 0, "usage": global_stats})
    try: c, r = psutil.cpu_percent(0.1), psutil.virtual_memory().percent
    except: c, r = 0, 0
    return jsonify({"cpu": c, "ram": r, "usage": global_stats})

@app.route('/download-report')
def download_report():
    if not session.get('is_admin'):
        return "Unauthorized Access.", 401
    try:
        now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        report = f"System Report - {now}\nStats: {global_stats}"
        return Response(report, mimetype="text/plain", headers={"Content-disposition": "attachment; filename=System_Report.txt"})
    except Exception as e: return str(e), 500

# ==============================================================================
#                               FEATURE ROUTES
# ==============================================================================

@app.route('/chat', methods=['POST'])
def chat():
    increment_stat('chat_msgs')
    try:
        msg = request.form.get('message', '')
        if not msg: return jsonify({"success": False, "error": "Empty message"}), 400

        # 1. RETRIEVE HISTORY from Session
        history = session.get('chat_history', [])
        
        # 2. DEFINE SYSTEM PERSONA (Expert Assistant)
        system_instruction = {
            "role": "system", 
            "content": (
                "You are an advanced AI Enterprise Assistant. Your goal is to be accurate, professional, and helpful.\n"
                "Rules:\n"
                "1. Be direct and concise. Avoid fluff.\n"
                "2. If asking about code, provide efficient, commented code blocks.\n"
                "3. If you do not know an answer, admit it. Do not hallucinate facts.\n"
                "4. Use formatting (bolding, lists) to make answers readable."
            )
        }

        # 3. BUILD MESSAGE CHAIN
        messages_payload = [system_instruction] + history + [{"role": "user", "content": msg}]

        # 4. GET AI RESPONSE
        ai_response = get_chat_response(messages_payload)

        # 5. UPDATE HISTORY (Append current exchange)
        history.append({"role": "user", "content": msg})
        history.append({"role": "assistant", "content": ai_response})
        
        # Keep only the last 10 messages (5 turns) to save session space
        if len(history) > 10:
            history = history[-10:]

        # Save back to session
        session['chat_history'] = history

        return jsonify({"success": True, "response": ai_response})

    except Exception as e: 
        return jsonify({"success": False, "error": str(e)}), 500

# Route to clear chat memory manually
@app.route('/clear-chat', methods=['POST'])
def clear_chat():
    session.pop('chat_history', None)
    return jsonify({"success": True})

@app.route('/generate-minutes', methods=['POST'])
def generate_minutes():
    increment_stat('text_gen')
    try:
        notes = request.form.get('notes', '')
        res = get_safe_ai_response(f"Create structured Meeting Minutes based on these notes:\n{notes}")
        return jsonify({"success": True, "minutes": res if res else "Error"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/make-ppt', methods=['POST'])
def make_ppt():
    increment_stat('text_gen')
    try:
        topic = request.form.get('topic', '')
        src_text = request.form.get('source_text', '')
        template_file = request.files.get('template_file')
        
        if template_file and template_file.filename != '':
            temp_path = os.path.join(STATIC_FOLDER, f"temp_{uuid.uuid4().hex}.pptx")
            template_file.save(temp_path)
            prs = Presentation(temp_path)
        else:
            prs = Presentation()
            
        prompt = (f"Create a presentation about {topic}. {src_text}. "
                  "Format exactly like this:\nSLIDE_TITLE: [Title]\nBULLET: [Point 1]\nBULLET: [Point 2]")
        
        content = get_safe_ai_response(prompt)
        if not content: content = f"SLIDE_TITLE: {topic}\nBULLET: Content generation failed."
        
        lines = content.split('\n')
        curr = None
        for line in lines:
            clean = line.strip().replace('*', '').replace('#', '')
            if "SLIDE_TITLE:" in clean:
                try:
                    layout_index = 1 if len(prs.slide_layouts) > 1 else 0
                    curr = prs.slides.add_slide(prs.slide_layouts[layout_index])
                    curr.shapes.title.text = clean.split("SLIDE_TITLE:", 1)[1].strip()
                    curr.placeholders[1].text_frame.word_wrap = True
                except: pass
            elif "BULLET:" in clean and curr:
                try:
                    p = curr.placeholders[1].text_frame.add_paragraph()
                    p.text = clean.split("BULLET:", 1)[1].strip()
                    p.level = 0
                except: pass
                
        fname = f"ppt_{uuid.uuid4().hex[:10]}.pptx"
        prs.save(os.path.join(STATIC_FOLDER, fname))
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

# --- UPDATED: TEXT TO AUDIO (Fixed Logic) ---
@app.route('/text-to-audio', methods=['POST'])
def text_to_audio():
    increment_stat('audio_gen')
    try:
        text = request.form.get('text', '')
        # Get language code directly (e.g., 'hi', 'ta', 'en')
        target_lang = request.form.get('target_language', 'en').strip()
        
        # Fallback: if empty or too long, default to english
        if not target_lang or len(target_lang) > 5:
            target_lang = 'en'
            
        fname = f"audio_{uuid.uuid4().hex[:10]}.mp3"
        
        # Generate Audio using gTTS
        tts = gTTS(text=text, lang=target_lang, slow=False)
        tts.save(os.path.join(STATIC_FOLDER, fname))

        return jsonify({
            "success": True, 
            "file_url": f"/static/{fname}", 
            "translated_text": text 
        })
    except Exception as e: 
        print(f"TTS Error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/audio-to-text', methods=['POST'])
def audio_to_text():
    increment_stat('transcribe')
    try:
        if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
        file = request.files['file']
        language_code = request.form.get('language', 'en-US') 
        if file.filename == '': return jsonify({"success": False, "error": "No file selected"}), 400

        filename = f"temp_rec_{uuid.uuid4().hex}.wav"
        filepath = os.path.join(STATIC_FOLDER, filename)
        file.save(filepath)

        recognizer = sr.Recognizer()
        with sr.AudioFile(filepath) as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data, language=language_code)

        if os.path.exists(filepath): os.remove(filepath)
        return jsonify({"success": True, "text": text})

    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/translate', methods=['POST'])
def translate():
    increment_stat('text_gen')
    try:
        text = request.form.get('text', '').strip()
        target_lang = request.form.get('target_language', 'en').strip()
        if not text: return jsonify({"success": False, "error": "Missing text"}), 400
        res = get_safe_ai_response(f"Translate this to {target_lang}: {text}")
        return jsonify({"success": True, "translation": res if res else "Error"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/generate-email', methods=['POST'])
def generate_email():
    increment_stat('text_gen')
    try:
        to, topic = request.form.get('recipient', ''), request.form.get('topic', '')
        res = get_safe_ai_response(f"Write a professional email to {to} about {topic}.")
        return jsonify({"success": True, "email_content": res if res else "Busy"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/text-to-pdf', methods=['POST'])
def text_to_pdf():
    increment_stat('pdf_gen')
    try:
        h = request.form.get('html_content', '')
        styled_html = f"""
        <html>
        <head><style>
            body {{ font-family: Helvetica; font-size: 12px; margin: 0; padding: 20px; }}
            p {{ margin-bottom: 5px; line-height: 1.4; }}
        </style></head>
        <body>{h}</body>
        </html>"""
        
        fname = f"doc_{uuid.uuid4().hex[:10]}.pdf"
        with open(os.path.join(STATIC_FOLDER, fname), "w+b") as f: 
            pisa.CreatePDF(BytesIO(styled_html.encode('utf-8')), dest=f)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

# --- UPDATED: QUIZ ROUTE (ABC Options) ---
@app.route('/generate-quiz', methods=['POST'])
def generate_quiz():
    increment_stat('quiz_gen')
    try:
        topic = request.form.get('topic', '')
        count = request.form.get('count', '5')
        
        # FIX: Prompt explicitly tells AI NOT to add A/B/C inside the text
        prompt = (f"Create a {count}-question Multiple Choice Quiz about: {topic}. "
                  f"Output strictly as HTML. "
                  f"Use <h3> for questions. "
                  f"Use <ol> (ordered list) for the options. "
                  f"IMPORTANT: Do NOT write 'A.', 'B.', 'C.' inside the <li> tags. Just write the answer text directly. "
                  f"At the end, add an <h2>Answer Key</h2> section as an HTML Table with columns 'Question' and 'Correct Answer'. "
                  f"Do not use markdown blocks.")
        
        res = get_safe_ai_response(prompt)
        
        if not res: return jsonify({"success": False, "error": "AI Service Busy"}), 500
        
        clean_html = res.replace("```html", "").replace("```", "").strip()
        
        # PDF CSS WRAPPER
        pdf_html = f"""
        <html>
            <head>
                <style>
                    @page {{ size: A4; margin: 1cm; }}
                    body {{ font-family: Helvetica, sans-serif; font-size: 11px; color: #000; }}
                    h1 {{ text-align: center; border-bottom: 1px solid #000; padding-bottom: 5px; margin-bottom: 15px; font-size: 18px; }}
                    h2 {{ color: #b30000; margin-top: 20px; margin-bottom: 10px; font-size: 14px; border-bottom: 1px solid #ccc; }}
                    h3 {{ color: #003366; margin-top: 10px; margin-bottom: 5px; font-size: 12px; }}
                    
                    /* This CSS adds the A, B, C, D automatically */
                    ol {{ margin-top: 0; margin-bottom: 10px; padding-left: 20px; list-style-type: upper-alpha; }}
                    li {{ margin-bottom: 2px; padding-left: 5px; }}
                    
                    table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
                    th, td {{ border: 1px solid #999; padding: 4px; text-align: left; font-size: 10px; }}
                    th {{ background-color: #eee; }}
                </style>
            </head>
            <body>
                <h1>Quiz: {topic}</h1>
                {clean_html}
            </body>
        </html>
        """
        
        fname = f"quiz_{uuid.uuid4().hex[:10]}.pdf"
        with open(os.path.join(STATIC_FOLDER, fname), "w+b") as f:
            pisa.CreatePDF(BytesIO(pdf_html.encode('utf-8')), dest=f)
            
        return jsonify({"success": True, "quiz": clean_html, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500
    
@app.route('/review-code', methods=['POST'])
def review_code():
    increment_stat('code_review')
    try:
        code = request.form.get('code', '')
        res = get_safe_ai_response(f"Review this code and suggest improvements:\n{code}")
        return jsonify({"success": True, "review": res if res else "Error"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/convert-file', methods=['POST'])
def convert_file():
    increment_stat('file_conv')
    try:
        if 'file' not in request.files: return jsonify({"success": False, "error": "No file uploaded"}), 400
        file = request.files['file']
        target_format = request.form.get('format', 'PNG').upper()
        if file.filename == '': return jsonify({"success": False, "error": "No file selected"}), 400
        
        img = PIL.Image.open(file)
        if target_format in ['JPEG', 'JPG', 'PDF']:
            img = img.convert('RGB')
            
        new_filename = f"converted_{uuid.uuid4().hex[:10]}.{target_format.lower()}"
        save_path = os.path.join(STATIC_FOLDER, new_filename)
        img.save(save_path, target_format if target_format != 'JPG' else 'JPEG')
        return jsonify({"success": True, "file_url": f"/static/{new_filename}"})
    except Exception as e: return jsonify({"success": False, "error": f"Error: {str(e)}"}), 500

@app.route('/compress-image', methods=['POST'])
def compress_image():
    increment_stat('compression')
    try:
        if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
        file = request.files['file']
        target_kb = int(request.form.get('target_kb', '500'))
        
        img = PIL.Image.open(file)
        img = img.convert('RGB')
        
        output_io = BytesIO()
        img.save(output_io, format='JPEG', quality=30, optimize=True)
            
        new_filename = f"compressed_{uuid.uuid4().hex[:10]}.jpg"
        save_path = os.path.join(STATIC_FOLDER, new_filename)
        with open(save_path, "wb") as f:
            f.write(output_io.getvalue())
        return jsonify({"success": True, "file_url": f"/static/{new_filename}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/video-to-audio', methods=['POST'])
def video_to_audio():
    increment_stat('vid_audio')
    try:
        if 'file' not in request.files: return jsonify({"success": False, "error": "No video file"}), 400
        file = request.files['file']
        
        temp_vid_name = f"temp_vid_{uuid.uuid4().hex[:10]}.mp4"
        temp_vid_path = os.path.join(STATIC_FOLDER, temp_vid_name)
        file.save(temp_vid_path)
        
        audio_name = f"extracted_{uuid.uuid4().hex[:10]}.mp3"
        audio_path = os.path.join(STATIC_FOLDER, audio_name)
        
        clip = VideoFileClip(temp_vid_path)
        clip.audio.write_audiofile(audio_path, logger=None)
        clip.close()
        
        if os.path.exists(temp_vid_path): os.remove(temp_vid_path)
        return jsonify({"success": True, "file_url": f"/static/{audio_name}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=True)